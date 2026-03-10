import io
import os
import re
import json
import difflib
from typing import Dict, Tuple, Any

import duckdb
import numpy as np
import pandas as pd
import requests
import streamlit as st
from dateutil import parser
from pptx import Presentation
from pptx.util import Inches, Pt


# -----------------------------
# Page setup
# -----------------------------
st.set_page_config(
    page_title="AI Initiative Discovery Agent",
    page_icon="📊",
    layout="wide",
)

st.title("📊 AI Initiative Discovery Agent")
st.caption(
    "Upload datasets → clean & normalize → analyze → call n8n backend → generate executive insights → create PowerPoint slides"
)


# -----------------------------
# Constants
# -----------------------------
STANDARD_VALUE_MAP = {
    "status": {
        "close": "closed",
        "closed": "closed",
        "clsed": "closed",
        "complete": "closed",
        "completed": "closed",
        "resolved": "closed",
        "resolve": "closed",
        "done": "closed",
        "open": "open",
        "opened": "open",
        "reopen": "open",
        "reopened": "open",
        "new": "open",
        "in progress": "in_progress",
        "in-progress": "in_progress",
        "inprogress": "in_progress",
        "progressing": "in_progress",
        "ongoing": "in_progress",
        "wip": "in_progress",
        "working": "in_progress",
        "cancel": "cancelled",
        "cancelled": "cancelled",
        "canceled": "cancelled",
        "pending": "pending",
        "pendng": "pending",
        "awaiting": "pending",
        "on hold": "on_hold",
        "hold": "on_hold",
    },
    "priority": {
        "hi": "high",
        "high": "high",
        "hgh": "high",
        "urgent": "high",
        "critical": "high",
        "med": "medium",
        "medium": "medium",
        "mid": "medium",
        "normal": "medium",
        "low": "low",
        "lo": "low",
    },
    "yes_no": {
        "y": "yes",
        "yes": "yes",
        "true": "yes",
        "1": "yes",
        "n": "no",
        "no": "no",
        "false": "no",
        "0": "no",
    },
}

IMPORTANT_COLUMN_HINTS = {
    "status": "status",
    "ticket_status": "status",
    "case_status": "status",
    "sr_status": "status",
    "request_status": "status",
    "priority": "priority",
    "severity": "priority",
    "urgency": "priority",
    "churn_flag": "yes_no",
    "active_flag": "yes_no",
    "retained_flag": "yes_no",
    "yn": "yes_no",
}

DEFAULT_BUSINESS_GOAL = (
    "Identify customer pain points, derive initiative opportunities, recommend KPI/KRIs, "
    "and reduce analyst time spent on manual data cleaning, analysis, and brainstorming."
)


# -----------------------------
# Session state defaults
# -----------------------------
for key, value in {
    "analysis_complete": False,
    "analysis_output": "",
    "ppt_bytes": None,
    "generated_prompt": "",
    "cleaned_datasets": {},
    "cleaning_reports": {},
    "dataset_profiles": {},
    "dataset_category_insights": {},
    "dataset_numeric_summaries": {},
    "machine_findings": {},
    "join_key_report": {},
    "executive_json": None,
}.items():
    if key not in st.session_state:
        st.session_state[key] = value


# -----------------------------
# Helper functions
# -----------------------------
def standardize_column_name(col: Any) -> str:
    col = str(col).strip().lower()
    col = re.sub(r"[^\w\s]", "", col)
    col = re.sub(r"\s+", "_", col)
    return col


def clean_text_value(x: Any):
    if pd.isna(x):
        return x
    if isinstance(x, str):
        x = x.strip()
        x = re.sub(r"\s+", " ", x)
        if x.lower() in {"n/a", "na", "null", "none", "nil", "blank", "missing", "nan", ""}:
            return np.nan
    return x


def normalize_for_match(x: Any):
    if pd.isna(x):
        return x
    x = str(x).strip().lower()
    x = re.sub(r"[_\-]+", " ", x)
    x = re.sub(r"\s+", " ", x)
    return x


def detect_column_semantic_type(col_name: str):
    col = col_name.lower()
    if col in IMPORTANT_COLUMN_HINTS:
        return IMPORTANT_COLUMN_HINTS[col]
    if any(k in col for k in ["status", "ticket_status", "case_status", "sr_status", "request_status"]):
        return "status"
    if any(k in col for k in ["priority", "severity", "urgency"]):
        return "priority"
    if any(k in col for k in ["active", "flag", "indicator", "yn", "yes_no", "churned", "retained"]):
        return "yes_no"
    return None


def standardize_categorical_values(df: pd.DataFrame, max_unique_ratio: float = 0.2, fuzzy_threshold: float = 0.88):
    df = df.copy()
    standardization_report = {}

    for col in df.columns:
        if not (df[col].dtype == "object" or str(df[col].dtype) == "category"):
            continue

        non_null = df[col].dropna()
        if len(non_null) == 0:
            continue

        unique_values = non_null.astype(str).nunique()
        unique_ratio = unique_values / max(len(non_null), 1)
        if unique_ratio > max_unique_ratio:
            continue

        semantic_type = detect_column_semantic_type(col)
        original_series = df[col].copy()
        normalized = df[col].map(normalize_for_match)

        if semantic_type in STANDARD_VALUE_MAP:
            value_map = STANDARD_VALUE_MAP[semantic_type]
            df[col] = normalized.map(lambda x: value_map.get(x, x) if pd.notna(x) else x)
        else:
            df[col] = normalized

        value_counts = df[col].dropna().astype(str).value_counts()
        unique_cleaned = value_counts.index.tolist()

        canonical_values = []
        fuzzy_map = {}
        for val in unique_cleaned:
            if not canonical_values:
                canonical_values.append(val)
                fuzzy_map[val] = val
                continue
            best_match = difflib.get_close_matches(val, canonical_values, n=1, cutoff=fuzzy_threshold)
            if best_match:
                fuzzy_map[val] = best_match[0]
            else:
                canonical_values.append(val)
                fuzzy_map[val] = val

        df[col] = df[col].map(lambda x: fuzzy_map.get(str(x), x) if pd.notna(x) else x)

        changed_count = int(
            (original_series.astype(str).fillna("<<NA>>") != df[col].astype(str).fillna("<<NA>>")).sum()
        )
        if changed_count > 0:
            sample_before_after = pd.DataFrame(
                {
                    "before": original_series.astype(str).fillna("<<NA>>"),
                    "after": df[col].astype(str).fillna("<<NA>>"),
                }
            )
            sample_before_after = sample_before_after[
                sample_before_after["before"] != sample_before_after["after"]
            ].drop_duplicates().head(15)

            standardization_report[col] = {
                "semantic_type": semantic_type if semantic_type else "generic_categorical",
                "unique_before": int(original_series.dropna().astype(str).nunique()),
                "unique_after": int(df[col].dropna().astype(str).nunique()),
                "changed_rows": changed_count,
                "sample_mappings": dict(zip(sample_before_after["before"], sample_before_after["after"])),
            }

    return df, standardization_report


def try_parse_dates(df: pd.DataFrame, threshold: float = 0.7):
    df = df.copy()
    converted_cols = []
    for col in df.columns:
        if df[col].dtype == "object":
            sample = df[col].dropna().astype(str).head(50)
            if len(sample) == 0:
                continue
            parsed_success = 0
            for val in sample:
                try:
                    parser.parse(val)
                    parsed_success += 1
                except Exception:
                    pass
            if (parsed_success / len(sample)) >= threshold:
                try:
                    df[col] = pd.to_datetime(df[col], errors="coerce")
                    converted_cols.append(col)
                except Exception:
                    pass
    return df, converted_cols


def try_parse_numeric(df: pd.DataFrame, threshold: float = 0.8):
    df = df.copy()
    converted_cols = []
    for col in df.columns:
        if df[col].dtype == "object":
            series = df[col].dropna().astype(str).str.replace(",", "", regex=False).str.strip()
            if len(series) == 0:
                continue
            numeric_series = pd.to_numeric(series, errors="coerce")
            if numeric_series.notna().mean() >= threshold:
                df[col] = pd.to_numeric(
                    df[col].astype(str).str.replace(",", "", regex=False), errors="coerce"
                )
                converted_cols.append(col)
    return df, converted_cols


def remove_duplicates(df: pd.DataFrame):
    before = len(df)
    df = df.drop_duplicates()
    after = len(df)
    return df, before - after


def handle_missing_values(df: pd.DataFrame):
    df = df.copy()
    fill_report = {}
    for col in df.columns:
        if pd.api.types.is_numeric_dtype(df[col]):
            missing_count = int(df[col].isna().sum())
            if missing_count > 0:
                median_val = df[col].median()
                df[col] = df[col].fillna(median_val)
                fill_report[col] = {
                    "strategy": "filled_numeric_with_median",
                    "missing_filled": missing_count,
                }
        elif pd.api.types.is_datetime64_any_dtype(df[col]):
            pass
        else:
            missing_count = int(df[col].isna().sum())
            if missing_count > 0:
                mode_vals = df[col].mode(dropna=True)
                fill_val = mode_vals.iloc[0] if len(mode_vals) > 0 else "unknown"
                df[col] = df[col].fillna(fill_val)
                fill_report[col] = {
                    "strategy": "filled_text_with_mode_or_unknown",
                    "missing_filled": missing_count,
                }
    return df, fill_report


def detect_outliers_iqr(df: pd.DataFrame):
    outlier_report = {}
    numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    for col in numeric_cols:
        series = df[col].dropna()
        if len(series) < 5:
            continue
        q1 = series.quantile(0.25)
        q3 = series.quantile(0.75)
        iqr = q3 - q1
        if iqr == 0:
            continue
        lower = q1 - 1.5 * iqr
        upper = q3 + 1.5 * iqr
        count = int(((df[col] < lower) | (df[col] > upper)).sum())
        outlier_report[col] = {
            "outlier_count": count,
            "lower_bound": float(lower),
            "upper_bound": float(upper),
        }
    return outlier_report


def clean_dataframe(df: pd.DataFrame, dataset_name: str = "dataset"):
    original_shape = df.shape
    df = df.copy()
    cleaning_report = {
        "dataset_name": dataset_name,
        "original_rows": int(original_shape[0]),
        "original_columns": int(original_shape[1]),
        "steps": {},
    }

    old_cols = df.columns.tolist()
    new_cols = [standardize_column_name(c) for c in df.columns]
    df.columns = new_cols
    cleaning_report["steps"]["column_standardization"] = {
        "old_columns": old_cols,
        "new_columns": new_cols,
    }

    df = df.apply(lambda col: col.map(clean_text_value))
    cleaning_report["steps"]["text_cleaning"] = "trimmed spaces, normalized blanks/null-like values"

    df, categorical_standardization_report = standardize_categorical_values(df)
    cleaning_report["steps"]["categorical_value_standardization"] = categorical_standardization_report

    df, duplicates_removed = remove_duplicates(df)
    cleaning_report["steps"]["duplicates_removed"] = duplicates_removed

    df, date_cols_converted = try_parse_dates(df)
    cleaning_report["steps"]["date_columns_converted"] = date_cols_converted

    df, numeric_cols_converted = try_parse_numeric(df)
    cleaning_report["steps"]["numeric_columns_converted"] = numeric_cols_converted

    df, missing_fill_report = handle_missing_values(df)
    cleaning_report["steps"]["missing_value_handling"] = missing_fill_report

    outlier_report = detect_outliers_iqr(df)
    cleaning_report["steps"]["outlier_detection"] = outlier_report

    cleaning_report["final_rows"] = int(df.shape[0])
    cleaning_report["final_columns"] = int(df.shape[1])
    return df, cleaning_report


def load_file(uploaded_file) -> pd.DataFrame | None:
    lower = uploaded_file.name.lower()
    if lower.endswith(".csv"):
        try:
            return pd.read_csv(uploaded_file)
        except UnicodeDecodeError:
            uploaded_file.seek(0)
            return pd.read_csv(uploaded_file, encoding="latin1")
    if lower.endswith(".xlsx") or lower.endswith(".xls"):
        return pd.read_excel(uploaded_file)
    return None


def profile_dataframe(df: pd.DataFrame):
    return {
        "rows": int(df.shape[0]),
        "columns": int(df.shape[1]),
        "column_names": df.columns.tolist(),
        "dtypes": {k: str(v) for k, v in df.dtypes.to_dict().items()},
        "missing_values_after_cleaning": df.isna().sum().sort_values(ascending=False).head(10).to_dict(),
        "sample_rows": df.head(3).astype(str).to_dict(orient="records"),
    }


def get_top_categories(df: pd.DataFrame, max_cols: int = 5, top_n: int = 5):
    results = {}
    object_cols = df.select_dtypes(include=["object"]).columns.tolist()
    for col in object_cols[:max_cols]:
        try:
            results[col] = df[col].astype(str).value_counts(dropna=False).head(top_n).to_dict()
        except Exception:
            pass
    return results


def get_numeric_summary(df: pd.DataFrame):
    numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    if not numeric_cols:
        return {}
    return df[numeric_cols].describe().round(2).to_dict()


def detect_possible_join_keys(datasets_dict: Dict[str, pd.DataFrame]):
    all_columns = {}
    for dataset_name, df in datasets_dict.items():
        for col in df.columns:
            all_columns.setdefault(col, []).append(dataset_name)
    shared_columns = {col: ds_list for col, ds_list in all_columns.items() if len(ds_list) > 1}
    key_patterns = ["id", "account", "customer", "service", "number", "no"]
    likely_keys = {
        col: ds_list
        for col, ds_list in shared_columns.items()
        if any(k in col.lower() for k in key_patterns)
    }
    return {"shared_columns": shared_columns, "likely_join_keys": likely_keys}


def build_machine_findings(datasets: Dict[str, pd.DataFrame]):
    con = duckdb.connect()
    findings = {}

    for idx, (name, df) in enumerate(datasets.items()):
        safe_base = standardize_column_name(name) or "dataset"
        relation_name = f"ds_{idx}_{safe_base}"
        con.register(relation_name, df)

        item = {}
        try:
            item["row_count"] = con.execute(
                f'SELECT COUNT(*) AS cnt FROM "{relation_name}"'
            ).df().to_dict(orient="records")
        except Exception as e:
            item["row_count_error"] = str(e)

        text_cols = df.select_dtypes(include=["object"]).columns.tolist()
        if text_cols:
            top_col = text_cols[0]
            try:
                q = f'''
                SELECT "{top_col}" AS category, COUNT(*) AS total
                FROM "{relation_name}"
                GROUP BY 1
                ORDER BY total DESC
                LIMIT 10
                '''
                item["top_category_table"] = con.execute(q).df().to_dict(orient="records")
            except Exception as e:
                item["top_category_table_error"] = str(e)

        findings[name] = item

    return findings


def build_analysis_payload(
    business_goal: str,
    cleaning_reports: dict,
    dataset_profiles: dict,
    dataset_category_insights: dict,
    dataset_numeric_summaries: dict,
    machine_findings: dict,
    join_key_report: dict,
):
    return {
        "business_goal": business_goal,
        "cleaning_reports": cleaning_reports,
        "dataset_profiles": dataset_profiles,
        "dataset_category_insights": dataset_category_insights,
        "dataset_numeric_summaries": dataset_numeric_summaries,
        "machine_findings": machine_findings,
        "join_key_report": join_key_report,
    }


def build_analysis_prompt_from_payload(payload: dict):
    return f"""
You are a senior strategy consultant, customer experience expert, and data analyst.

Context:
The machine has already performed initial data handling and cleaning on all uploaded datasets, including:
- column name standardization
- text cleanup
- null normalization
- duplicate removal
- date parsing
- numeric parsing
- missing value treatment
- outlier detection summary
- category normalization for low-cardinality fields
- normalization of equivalent business values such as close/closed/clsed into one category where appropriate
- initial profiling and summary tables

Business goal:
{payload["business_goal"]}

Cleaning reports:
{json.dumps(payload["cleaning_reports"], indent=2, default=str)}

Dataset profiles after cleaning:
{json.dumps(payload["dataset_profiles"], indent=2, default=str)}

Category insights:
{json.dumps(payload["dataset_category_insights"], indent=2, default=str)}

Numeric summaries:
{json.dumps(payload["dataset_numeric_summaries"], indent=2, default=str)}

Machine-generated findings:
{json.dumps(payload["machine_findings"], indent=2, default=str)}

Possible join keys across datasets:
{json.dumps(payload["join_key_report"], indent=2, default=str)}

Instructions:
Return strict JSON only using this exact schema:
{{
  "executive_problem_statement": "...",
  "key_insights": ["...", "...", "..."],
  "data_quality_limitations": ["...", "..."],
  "root_cause_hypotheses": ["...", "...", "..."],
  "initiative_opportunities": [
    {{
      "initiative_name": "...",
      "issue_solved": "...",
      "why_it_matters": "...",
      "expected_business_value": "...",
      "effort_level": "low|medium|high"
    }}
  ],
  "kpi_recommendations": [
    {{
      "leading_kpi": "...",
      "lagging_kpi": "...",
      "suggested_baseline": "...",
      "suggested_target": "...",
      "why_it_matters": "..."
    }}
  ],
  "execution_plan_30_60_90": {{
    "days_0_30": ["...", "..."],
    "days_31_60": ["...", "..."],
    "days_61_90": ["...", "..."]
  }},
  "slide_ready_summary": {{
    "slide_title": "...",
    "subtitle": "...",
    "bullets": ["...", "...", "...", "..."],
    "suggested_chart": "...",
    "expected_business_impact": "...",
    "expected_productivity_gain": "..."
  }}
}}
""".strip()


def call_n8n_webhook(webhook_url: str, payload: dict) -> Tuple[dict | None, str]:
    try:
        response = requests.post(webhook_url, json=payload, timeout=180)
        response.raise_for_status()
        data = response.json()
        return data, json.dumps(data)
    except Exception as e:
        return None, f"n8n webhook error: {e}"


def safe_list(value):
    return value if isinstance(value, list) else []


def parse_fallback_json(raw_text: str):
    try:
        return json.loads(raw_text)
    except Exception:
        return None


def add_bullets(tf, bullets, level=0, font_size=20):
    tf.clear()
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = str(bullet)
        p.level = level
        p.font.size = Pt(font_size)


def _set_title_and_body(slide, title_text: str, bullets=None, body_text: str | None = None):
    title_set = False
    body_set = False

    if hasattr(slide.shapes, "title") and slide.shapes.title is not None:
        slide.shapes.title.text = title_text
        title_set = True

    for shape in slide.placeholders:
        if getattr(shape, "placeholder_format", None) is None:
            continue
        if not body_set and shape.has_text_frame and shape != getattr(slide.shapes, "title", None):
            tf = shape.text_frame
            if bullets is not None:
                add_bullets(tf, bullets, font_size=18)
            else:
                tf.clear()
                tf.paragraphs[0].text = body_text or ""
                tf.paragraphs[0].font.size = Pt(20)
            body_set = True

    if not title_set:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        p = tx.text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True

    if not body_set:
        tx = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.8), Inches(4.8))
        tf = tx.text_frame
        if bullets is not None:
            add_bullets(tf, bullets, font_size=18)
        else:
            tf.paragraphs[0].text = body_text or ""
            tf.paragraphs[0].font.size = Pt(20)


def create_pptx(executive_json: dict, app_title: str = "AI Initiative Discovery Agent") -> bytes:
    template_path = "template_exec_deck.pptx"
    prs = Presentation(template_path) if os.path.exists(template_path) else Presentation()

    required_slides = 7
    while len(prs.slides) < required_slides:
        prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])

    slide = prs.slides[0]
    _set_title_and_body(
        slide,
        executive_json.get("slide_ready_summary", {}).get("slide_title", app_title),
        body_text=executive_json.get("slide_ready_summary", {}).get(
            "subtitle", "Executive summary generated from uploaded datasets"
        ),
    )

    slide = prs.slides[1]
    _set_title_and_body(
        slide,
        "Executive Problem Statement",
        body_text=executive_json.get("executive_problem_statement", ""),
    )

    slide = prs.slides[2]
    _set_title_and_body(slide, "Key Insights from Data", bullets=safe_list(executive_json.get("key_insights")))

    slide = prs.slides[3]
    _set_title_and_body(
        slide, "Root Cause Hypotheses", bullets=safe_list(executive_json.get("root_cause_hypotheses"))
    )

    initiative_bullets = []
    for item in safe_list(executive_json.get("initiative_opportunities")):
        initiative_bullets.append(
            f"{item.get('initiative_name', 'Initiative')}: {item.get('issue_solved', '')} | "
            f"Value: {item.get('expected_business_value', '')} | Effort: {item.get('effort_level', '')}"
        )
    slide = prs.slides[4]
    _set_title_and_body(slide, "Initiative Opportunities", bullets=initiative_bullets)

    kpi_bullets = []
    for item in safe_list(executive_json.get("kpi_recommendations")):
        kpi_bullets.append(
            f"Leading KPI: {item.get('leading_kpi', '')} | Lagging KPI: {item.get('lagging_kpi', '')} | "
            f"Baseline: {item.get('suggested_baseline', '')} | Target: {item.get('suggested_target', '')}"
        )
    slide = prs.slides[5]
    _set_title_and_body(slide, "KPI Recommendations", bullets=kpi_bullets)

    plan = executive_json.get("execution_plan_30_60_90", {})
    summary = executive_json.get("slide_ready_summary", {})
    final_bullets = [
        "0-30 days: " + " | ".join(safe_list(plan.get("days_0_30"))),
        "31-60 days: " + " | ".join(safe_list(plan.get("days_31_60"))),
        "61-90 days: " + " | ".join(safe_list(plan.get("days_61_90"))),
    ] + safe_list(summary.get("bullets")) + [
        f"Suggested chart: {summary.get('suggested_chart', '')}",
        f"Expected business impact: {summary.get('expected_business_impact', '')}",
        f"Expected productivity gain: {summary.get('expected_productivity_gain', '')}",
    ]
    slide = prs.slides[6]
    _set_title_and_body(slide, "Execution Plan & Summary", bullets=final_bullets)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def render_executive_output(data: dict):
    st.subheader("Executive Problem Statement")
    st.write(data.get("executive_problem_statement", "-"))

    c1, c2 = st.columns(2)
    with c1:
        st.subheader("Key Insights")
        for item in safe_list(data.get("key_insights")):
            st.markdown(f"- {item}")
    with c2:
        st.subheader("Data Quality / Analysis Limitations")
        for item in safe_list(data.get("data_quality_limitations")):
            st.markdown(f"- {item}")

    c3, c4 = st.columns(2)
    with c3:
        st.subheader("Root Cause Hypotheses")
        for item in safe_list(data.get("root_cause_hypotheses")):
            st.markdown(f"- {item}")
    with c4:
        st.subheader("30-60-90 Plan")
        plan = data.get("execution_plan_30_60_90", {})
        st.markdown("**0-30 days**")
        for item in safe_list(plan.get("days_0_30")):
            st.markdown(f"- {item}")
        st.markdown("**31-60 days**")
        for item in safe_list(plan.get("days_31_60")):
            st.markdown(f"- {item}")
        st.markdown("**61-90 days**")
        for item in safe_list(plan.get("days_61_90")):
            st.markdown(f"- {item}")

    st.subheader("Initiative Opportunities")
    initiatives = pd.DataFrame(safe_list(data.get("initiative_opportunities")))
    if not initiatives.empty:
        st.dataframe(initiatives, use_container_width=True)

    st.subheader("KPI Recommendations")
    kpis = pd.DataFrame(safe_list(data.get("kpi_recommendations")))
    if not kpis.empty:
        st.dataframe(kpis, use_container_width=True)

    st.subheader("Slide-ready Summary")
    slide_ready = data.get("slide_ready_summary", {})
    st.markdown(f"**Title:** {slide_ready.get('slide_title', '-')}")
    st.markdown(f"**Subtitle:** {slide_ready.get('subtitle', '-')}")
    for bullet in safe_list(slide_ready.get("bullets")):
        st.markdown(f"- {bullet}")
    st.markdown(f"**Suggested chart:** {slide_ready.get('suggested_chart', '-')}")
    st.markdown(f"**Expected business impact:** {slide_ready.get('expected_business_impact', '-')}")
    st.markdown(f"**Expected productivity gain:** {slide_ready.get('expected_productivity_gain', '-')}")


# -----------------------------
# Sidebar controls
# -----------------------------
with st.sidebar:
    st.header("Settings")
    business_goal = st.text_area("Business goal", value=DEFAULT_BUSINESS_GOAL, height=120)
    auto_call_n8n = st.toggle("Auto-call n8n backend", value=True)
    n8n_webhook_url = st.text_input(
        "n8n webhook URL",
        value=st.secrets.get("N8N_WEBHOOK_URL", "") if hasattr(st, "secrets") else "",
    )
    st.caption("For deployment, store the webhook in Streamlit secrets as N8N_WEBHOOK_URL.")
    st.caption("Optional: add template_exec_deck.pptx to the repo root for branded slides.")


# -----------------------------
# Main UI tabs
# -----------------------------
tab1, tab2, tab3, tab4 = st.tabs([
    "1. Upload & Clean",
    "2. Executive Insights",
    "3. Slides",
    "4. Presenter Guide",
])

with tab1:
    uploaded_files = st.file_uploader(
        "Upload CSV or Excel files",
        type=["csv", "xlsx", "xls"],
        accept_multiple_files=True,
        help="Upload one or more datasets.",
    )

    run_analysis = st.button("Run Analysis", type="primary", use_container_width=True)

    if run_analysis:
        if not uploaded_files:
            st.error("Please upload at least one file.")
        else:
            raw_datasets = {}
            cleaned_datasets = {}
            cleaning_reports = {}

            with st.spinner("Loading, cleaning, and profiling datasets..."):
                for file in uploaded_files:
                    df = load_file(file)
                    if df is None:
                        st.warning(f"Skipped unsupported file: {file.name}")
                        continue

                    dataset_name = re.sub(r"\.[^.]+$", "", file.name)
                    raw_datasets[dataset_name] = df

                    cleaned_df, report = clean_dataframe(df, dataset_name=dataset_name)
                    cleaned_datasets[dataset_name] = cleaned_df
                    cleaning_reports[dataset_name] = report

                if not cleaned_datasets:
                    st.error("No supported files were loaded.")
                else:
                    dataset_profiles = {name: profile_dataframe(df) for name, df in cleaned_datasets.items()}
                    dataset_category_insights = {name: get_top_categories(df) for name, df in cleaned_datasets.items()}
                    dataset_numeric_summaries = {name: get_numeric_summary(df) for name, df in cleaned_datasets.items()}
                    join_key_report = detect_possible_join_keys(cleaned_datasets)
                    machine_findings = build_machine_findings(cleaned_datasets)

                    payload = build_analysis_payload(
                        business_goal,
                        cleaning_reports,
                        dataset_profiles,
                        dataset_category_insights,
                        dataset_numeric_summaries,
                        machine_findings,
                        join_key_report,
                    )

                    prompt = build_analysis_prompt_from_payload(payload)

                    st.session_state.cleaned_datasets = cleaned_datasets
                    st.session_state.cleaning_reports = cleaning_reports
                    st.session_state.dataset_profiles = dataset_profiles
                    st.session_state.dataset_category_insights = dataset_category_insights
                    st.session_state.dataset_numeric_summaries = dataset_numeric_summaries
                    st.session_state.join_key_report = join_key_report
                    st.session_state.machine_findings = machine_findings
                    st.session_state.generated_prompt = prompt
                    st.session_state.analysis_complete = True
                    st.session_state.analysis_output = ""
                    st.session_state.executive_json = None
                    st.session_state.ppt_bytes = None

                    if auto_call_n8n and n8n_webhook_url:
                        with st.spinner("Calling n8n backend..."):
                            parsed_json, raw_text = call_n8n_webhook(n8n_webhook_url, payload)

                        if parsed_json is not None:
                            st.session_state.executive_json = parsed_json
                            st.session_state.analysis_output = raw_text
                            st.session_state.ppt_bytes = create_pptx(parsed_json)
                        else:
                            st.session_state.analysis_output = raw_text

                    st.success("Analysis completed.")

    if st.session_state.analysis_complete:
        st.subheader("Cleaned Datasets")
        for name, df in st.session_state.cleaned_datasets.items():
            with st.expander(f"{name} — {df.shape[0]} rows × {df.shape[1]} columns", expanded=False):
                st.dataframe(df.head(20), use_container_width=True)
                obj_cols = df.select_dtypes(include=["object"]).columns.tolist()
                if obj_cols:
                    top_col = obj_cols[0]
                    chart_df = df[top_col].astype(str).value_counts().head(10)
                    st.caption(f"Top values in {top_col}")
                    st.bar_chart(chart_df)

        st.subheader("Likely Join Keys")
        st.json(st.session_state.join_key_report)

        st.subheader("Cleaning Reports")
        st.json(st.session_state.cleaning_reports)

with tab2:
    if not st.session_state.analysis_complete:
        st.info("Run the analysis first in Tab 1.")
    else:
        if st.session_state.executive_json is not None:
            render_executive_output(st.session_state.executive_json)
        else:
            st.warning("Executive insights are not generated yet.")
            st.markdown("**Option A: Auto mode** — enter a working n8n webhook URL in the sidebar and rerun analysis.")
            st.markdown("**Option B: Manual mode** — copy the prompt below into any tool, then paste the JSON result back here.")

            st.subheader("Prompt")
            st.code(st.session_state.generated_prompt, language="text")

            pasted_json = st.text_area(
                "Paste JSON output here",
                height=320,
                placeholder="Paste the model or mock JSON response here...",
            )
            if st.button("Use Pasted JSON", use_container_width=True):
                parsed = parse_fallback_json(pasted_json)
                if parsed is None:
                    st.error("That is not valid JSON. Paste the exact JSON output.")
                else:
                    st.session_state.executive_json = parsed
                    st.session_state.analysis_output = pasted_json
                    st.session_state.ppt_bytes = create_pptx(parsed)
                    st.success("Executive insights loaded successfully.")
                    st.rerun()

with tab3:
    if st.session_state.executive_json is None:
        st.info("Generate executive insights first in Tab 2.")
    else:
        st.success("Slides generated.")
        slide_ready = st.session_state.executive_json.get("slide_ready_summary", {})
        st.markdown(f"**Deck title:** {slide_ready.get('slide_title', 'Executive Insights Deck')}")
        st.markdown(f"**Subtitle:** {slide_ready.get('subtitle', '-')}")

        if st.session_state.ppt_bytes is not None:
            st.download_button(
                label="Download PowerPoint (.pptx)",
                data=st.session_state.ppt_bytes,
                file_name="ai_initiative_executive_deck.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )

        st.subheader("Raw Executive JSON")
        st.json(st.session_state.executive_json)

with tab4:
    st.markdown(
        """
### Presenter flow
1. Open the Streamlit link.
2. Upload the datasets.
3. Click **Run Analysis**.
4. Open **Executive Insights**.
5. Open **Slides** and click **Download PowerPoint**.

### Demo backend mode
- The app calls an **n8n webhook**.
- n8n returns a **predefined JSON response**.
- This makes the demo stable and avoids live model/API quota issues.

### Template file
- File name: `template_exec_deck.pptx`
- Place it in the same folder as `app.py`.
- The app uses the first 7 slides of the template.
- Each slide should have a title placeholder and either a body placeholder or a text box area for content.

### Recommended requirements.txt
```txt
streamlit
pandas
duckdb
openpyxl
numpy
python-dateutil
python-pptx
requests
